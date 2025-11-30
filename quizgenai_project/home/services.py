import fitz
import google.generativeai as genai
import json
import requests
import olefile
import struct
import re
import os
from bs4 import BeautifulSoup
from django.conf import settings
from urllib.parse import urlparse, parse_qs
from pptx import Presentation

try:
    genai.configure(api_key=settings.GEMINI_API_KEY)
    model = genai.GenerativeModel('gemini-2.5-flash')
except Exception as e:
    print(f"Error configuring Gemini API in services.py: {e}")
    model = None

def extract_text_from_url(url):
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style", "nav", "footer", "header", "aside"]):
            script.decompose()
            
        text = soup.get_text(separator=' ', strip=True)
        return text
    except Exception as e:
        print(f"Error extracting text from URL: {e}")
        raise Exception(f"Could not extract text from the URL: {e}")

def generate_quiz_from_text(text, num_questions, custom_instructions):
    if model is None:
        raise Exception("Gemini API model is not configured.")

    if not text.strip():
        raise Exception("Could not extract any meaningful text.")

    prompt_base = f"""
    Based on the following text, create a multiple-choice quiz with {num_questions} questions.
    The response MUST be a valid JSON array and nothing else. Do not include any text, code block markers like ```json, or any other formatting before or after the JSON array.

    Each object in the array must have these exact keys:
    1. "question": A string for the question text.
    2. "options": An array of 4 strings representing the possible answers.
    3. "correctAnswer": The 0-based index of the correct answer within the "options" array.
    4. "explanation": A short explanation (1-2 sentences) of why the correct answer is correct.

    Example format:
    [
      {{"question": "What is the primary topic?", "options": ["A", "B", "C", "D"], "correctAnswer": 2, "explanation": "Option C is correct because..."}}
    ]
    IMPORTANT: You must generate questions that test the conceptual and scientific understanding of the topics presented in the text.
    GOOD QUESTIONS are about: chemical principles, scientific concepts, reactions, definitions, properties, and applications (e.g., "What is a thermodynamic function?", "What is the principle of potentiometry?").
    BAD QUESTIONS (DO NOT ASK): Trivial questions about the document's structure, layout, or metadata. This includes questions about page numbers, section headings, experiment numbers (e.g., "What is Experiment 1?"), assessment procedures, rubrics, or lists of content.
    """

    if custom_instructions and custom_instructions.strip():
        prompt_instructions = f"\n\nFollow these additional instructions when generating the quiz: {custom_instructions.strip()}\n"
    else:
        prompt_instructions = ""

    max_chars = 100000 
    if len(text) > max_chars:
        text = text[:max_chars] + "...(truncated)"

    prompt_text_section = f"""
    Here is the text to analyze:
    ---
    {text}
    ---
    """

    final_prompt = prompt_base + prompt_instructions + prompt_text_section

    try:
        response = model.generate_content(final_prompt)
    except Exception as api_error:
        print(f"Error calling Gemini API: {api_error}")
        raise Exception(f"Failed to communicate with the AI model: {api_error}")

    try:
        if not response.text:
            print(f"Gemini Response Feedback: {response.prompt_feedback}")
            raise Exception("The AI returned an empty response. This may be due to safety filters or the input content.")

        response_text = response.text.strip()
        
        # Robust JSON extraction: find the first '[' and the last ']'
        start_idx = response_text.find("[")
        end_idx = response_text.rfind("]")

        if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
            response_text = response_text[start_idx:end_idx+1]
        else:
            # Fallback: basic markdown cleanup if brackets aren't found correctly
            if response_text.startswith("```json"):
                response_text = response_text[7:]
            elif response_text.startswith("```"):
                response_text = response_text[3:]
            if response_text.endswith("```"):
                response_text = response_text[:-3]
            response_text = response_text.strip()

        quiz_data = json.loads(response_text)
        
        if not isinstance(quiz_data, list):
            raise Exception("The AI response is not a valid array of questions.")
        
        if len(quiz_data) == 0:
            raise Exception("The AI generated no questions. Please try again with a different PDF or fewer questions.")
        
        for i, question in enumerate(quiz_data):
            if not isinstance(question, dict):
                raise Exception(f"Question {i+1} is not a valid object.")
            if "question" not in question or "options" not in question or "correctAnswer" not in question:
                raise Exception(f"Question {i+1} is missing required fields (question, options, correctAnswer).")
            if not isinstance(question.get("options"), list) or len(question.get("options", [])) != 4:
                raise Exception(f"Question {i+1} does not have exactly 4 options.")
            # Explanation is optional for backward compatibility, but we expect it for new quizzes
            if "explanation" not in question:
                question["explanation"] = "No explanation provided."

    except json.JSONDecodeError as json_err:
        print(f"Failed to decode JSON. Raw response: {response.text}")
        raise Exception("The AI returned a malformed response. This sometimes happens with very large question counts. Try reducing the number of questions (max 25 recommended) or try again.")
    except Exception as parse_error:
        print(f"Error processing response: {parse_error}")
        if "JSON" not in str(parse_error):
            raise parse_error
        raise Exception("An error occurred while processing the AI response. Please try again with fewer questions.")

    return quiz_data

def generate_quiz_from_pdf(pdf_file, num_questions, custom_instructions):
    try:
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        pdf_text = "".join(page.get_text() for page in pdf_document)
    except Exception as fitz_error:
        print(f"Error opening or reading PDF: {fitz_error}")
        raise Exception(f"Could not process the PDF file: {fitz_error}")

    return generate_quiz_from_text(pdf_text, num_questions, custom_instructions)

def extract_text_from_ppt_legacy(ppt_file):
    try:
        if hasattr(ppt_file, 'seek'):
            ppt_file.seek(0)
        
        ole = olefile.OleFileIO(ppt_file)
        
        if not ole.exists('PowerPoint Document'):
            return ""
            
        stream = ole.openstream('PowerPoint Document')
        data = stream.read()
        
        text_runs = []
        idx = 0
        while idx < len(data):
            if idx + 8 > len(data):
                break
                
            rec_ver_inst = struct.unpack('<H', data[idx:idx+2])[0]
            rec_type = struct.unpack('<H', data[idx+2:idx+4])[0]
            rec_len = struct.unpack('<I', data[idx+4:idx+8])[0]
            
            rec_ver = rec_ver_inst & 0x000F
            
            if rec_type == 4000:
                if idx + 8 + rec_len <= len(data):
                    text_bytes = data[idx+8:idx+8+rec_len]
                    try:
                        text = text_bytes.decode('utf-16le')
                        text_runs.append(text)
                    except:
                        pass
                idx += 8 + rec_len
            elif rec_type == 4008:
                if idx + 8 + rec_len <= len(data):
                    text_bytes = data[idx+8:idx+8+rec_len]
                    try:
                        text = text_bytes.decode('latin-1')
                        text_runs.append(text)
                    except:
                        pass
                idx += 8 + rec_len
            elif rec_ver == 0x0F:
                idx += 8
            else:
                idx += 8 + rec_len
            
        return "\n".join(text_runs)
        
    except Exception as e:
        print(f"Error extracting text from legacy PPT: {e}")
        raise Exception(f"Could not extract text from .ppt file: {e}")

def generate_quiz_from_ppt(ppt_file, num_questions, custom_instructions):
    try:
        if ppt_file.name.lower().endswith('.ppt'):
             ppt_text = extract_text_from_ppt_legacy(ppt_file)
        elif ppt_file.name.lower().endswith('.pptx'):
            prs = Presentation(ppt_file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            ppt_text = "\n".join(text_runs)
        else:
             raise Exception("Unsupported file format. Please upload .ppt or .pptx.")
             
    except Exception as ppt_error:
        print(f"Error opening or reading PPT: {ppt_error}")
        raise Exception(f"Could not process the PPT file: {ppt_error}")

    return generate_quiz_from_text(ppt_text, num_questions, custom_instructions)

def generate_quiz_from_url(url, num_questions, custom_instructions):
    text = extract_text_from_url(url)
    return generate_quiz_from_text(text, num_questions, custom_instructions)


def estimate_generation_time(pdf_text_length, num_questions):
    base_min = 15
    base_max = 25
    
    size_factor_min = (pdf_text_length / 10000) * 5
    size_factor_max = (pdf_text_length / 10000) * 8
    
    questions_factor_min = int(num_questions) * 3
    questions_factor_max = int(num_questions) * 5
    
    total_min = int(base_min + size_factor_min + questions_factor_min)
    total_max = int(base_max + size_factor_max + questions_factor_max)
    
    total_min = min(total_min, 180)
    total_max = min(total_max, 180)
    
    return (total_min, total_max)

